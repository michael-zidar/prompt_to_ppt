from pathlib import Path
from typing import Dict, Any

from pptx import Presentation


def extract_style(pptx_path: Path) -> Dict[str, Any]:
    """Extract basic style information from a pptx file."""
    prs = Presentation(pptx_path)
    style = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "fonts": {},
        "theme_colors": {},
        "layouts": [],
    }

    if prs.slide_master and prs.slide_master.placeholders:
        for ph in prs.slide_master.placeholders:
            if ph.placeholder_format.type == 1:  # title
                style["fonts"]["title"] = ph.text_frame.paragraphs[0].font.name
            elif ph.placeholder_format.type == 2:  # body
                style["fonts"]["body"] = ph.text_frame.paragraphs[0].font.name

    for layout in prs.slide_layouts:
        placeholders = [ph.name for ph in layout.placeholders]
        style["layouts"].append({"name": layout.name, "placeholders": placeholders})

    return style
