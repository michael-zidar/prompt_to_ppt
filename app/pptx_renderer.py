from pathlib import Path
from typing import Dict, Any, List

from pptx import Presentation
from pptx.util import Inches


def build_pptx(slides: List[Dict[str, Any]], style: Dict[str, Any], out_path: Path) -> Path:
    prs = Presentation()
    for slide_data in slides:
        layout_name = slide_data.get("layout", "Title and Content")
        layout = prs.slide_layouts[0]
        for sl in prs.slide_layouts:
            if sl.name == layout_name:
                layout = sl
                break
        slide = prs.slides.add_slide(layout)
        elements = slide_data.get("elements", {})
        for placeholder in slide.placeholders:
            key = placeholder.name.lower()
            text = elements.get(key)
            if text:
                placeholder.text = text if isinstance(text, str) else "\n".join(text)
    prs.save(out_path)
    return out_path
